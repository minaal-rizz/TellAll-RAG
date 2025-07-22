# app/utils/extractors.py

import os
import fitz
import pytesseract
import shutil
import pandas as pd
import tempfile
from pdf2image import convert_from_path
from pptx import Presentation
from docx2pdf import convert


def extract_text(file_path):
    """
    Extract readable and page-aware text from supported document formats (PDF, DOCX, PPTX, XLSX).

    Args:
        file_path (str): Absolute path to the uploaded document.

    Returns:
        List[str]: List of strings where each item corresponds to a page/slide/row's text content.
    """
    try:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            doc = fitz.open(file_path)
            text = [p.get_text() for p in doc if p.get_text().strip()]
            if text:
                return text
            return [pytesseract.image_to_string(img) for img in convert_from_path(file_path)]

        elif ext == ".docx":
            with tempfile.TemporaryDirectory() as tmpdir:
                tmp_pdf_path = os.path.join(tmpdir, "converted.pdf")
                shutil.copy(file_path, os.path.join(tmpdir, "upload.docx"))
                convert(os.path.join(tmpdir, "upload.docx"), tmp_pdf_path)

                doc = fitz.open(tmp_pdf_path)
                text = [p.get_text() for p in doc if p.get_text().strip()]
                if text:
                    return text
                return [pytesseract.image_to_string(img) for img in convert_from_path(tmp_pdf_path)]

        elif ext == ".pptx":
            prs = Presentation(file_path)
            return [
                f"Slide {i+1}:\n" + "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
                ) for i, slide in enumerate(prs.slides)
            ]

        elif ext in [".xls", ".xlsx"]:
            text_chunks = []
            dfs = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in dfs.items():
                df = df.fillna("").astype(str)
                for i, row in df.iterrows():
                    row_text = ", ".join(row.tolist())
                    if row_text.strip():
                        text_chunks.append(f"Sheet: {sheet_name} | Row: {i+1}\n{row_text}")
            return text_chunks

    except Exception as e:
        print(f"❌ Error reading {file_path}: {e}")

    return []
