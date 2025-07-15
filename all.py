

import os
import sys
from datetime import datetime, timezone
from functools import lru_cache
import contextlib
from sentence_transformers import CrossEncoder #for re rankers
import gradio as gr
from pymongo import MongoClient
from pinecone import Pinecone
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
from docx import Document
import fitz
from pdf2image import convert_from_path
import pytesseract
import pandas as pd
import hashlib
import openai
import certifi
from groq import Groq

import tempfile
import shutil
from docx2pdf import convert


from dotenv import load_dotenv
load_dotenv()
NAMESPACE = os.getenv("NAMESPACE", "__default__")
print(f"✅ ENV LOADED: NAMESPACE={NAMESPACE}")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX = os.getenv("PINECONE_INDEX", "rag-index")
MONGO_URI = os.getenv("MONGO_URI")
print(os.getenv("MONGO_URI"))
print(f"📂 Using Pinecone namespace: '{NAMESPACE}' on index: '{PINECONE_INDEX}'")

# Import and initialize Groq client
openai.api_base = "https://api.groq.com/openai/v1"
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))


# Connect to MongoDB
try:
    client = MongoClient(MONGO_URI, tls=True,tlsCAFile=certifi.where())
    db = client["tellall"]
    docs_collection = db["documents"]
except Exception as e:
    print(f"❌ MongoDB connection failed: {e}")
    docs_collection = None

# Set up Pinecone and embedding model
pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(PINECONE_INDEX)
model = SentenceTransformer("all-MiniLM-L6-v2")
splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)

@lru_cache(maxsize=1)
def load_clients():
    """Load and cache the embedding model and Pinecone index."""
    return model, index

#sync mongo n pinecone 
def sync_mongo_pinecone():
    """Automatically sync Pinecone and MongoDB at startup using chunk-level IDs."""
    print("🔄 Syncing MongoDB and Pinecone...")

    if docs_collection is None:
        print("❌ MongoDB not connected.")
        return

    try:
        # Step 1: Collect all chunk IDs from MongoDB
        mongo_docs = list(docs_collection.find())
        mongo_chunk_ids = set()
        for doc in mongo_docs:
            chunk_ids = doc.get("chunk_ids", [])
            mongo_chunk_ids.update(chunk_ids)

        print(f"📦 Found {len(mongo_chunk_ids)} chunk IDs in MongoDB.")

        # Step 2: Query Pinecone for all stored chunks (up to 1000)
        try:
            res = index.query(
                vector=[0.0]*384,
                top_k=1000,
                include_metadata=True,
                namespace=NAMESPACE
            )
            pinecone_ids = {match["id"] for match in res.get("matches", [])}

            print(f"📡 Found {len(pinecone_ids)} chunk IDs in Pinecone.")

            # Step 3: Identify and delete orphaned Pinecone chunks
            orphaned = pinecone_ids - mongo_chunk_ids
            if orphaned:
                print(f"🧹 Deleting {len(orphaned)} stale Pinecone chunks...")
                index.delete(ids=list(orphaned), namespace=NAMESPACE)
            else:
                print("✅ No stale Pinecone chunks found.")

            print("✅ MongoDB ↔️ Pinecone sync complete.")
        except Exception as e:
            print(f"❌ Pinecone query failed: {e}")

    except Exception as e:
        print(f"❌ Sync failed: {e}")
sync_mongo_pinecone()  # Run sync at startup


'''def calculate_file_hash(filepath):
    """Calculate SHA-256 hash of the file to ensure uniqueness."""
    hasher = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hasher.update(chunk)
    return hasher.hexdigest()'''


def extract_text(file_path):
    """
Extract readable and page-aware text from supported document formats (PDF, DOCX, PPTX, XLSX).

Args:
    file_path (str): Absolute path to the uploaded document.

Returns:
    List[str]: List of strings where each item corresponds to a page/slide/row's text content.
"""

    try:
        ext = os.path.splitext(file_path)[1].lower()  # get the file extension in lowercase

        if ext == ".pdf":
            doc = fitz.open(file_path)  # open the pdf using pymupdf
            text = [p.get_text() for p in doc if p.get_text().strip()]  # extract text page by page
            if text:
                return text  # return extracted text if any
            return [pytesseract.image_to_string(img) for img in convert_from_path(file_path)]  # fallback to ocr

        elif ext == ".docx":
            with tempfile.TemporaryDirectory() as tmpdir:  # create a temporary directory
                tmp_pdf_path = os.path.join(tmpdir, "converted.pdf")  # define path for converted pdf
                shutil.copy(file_path, os.path.join(tmpdir, "upload.docx"))  # copy docx to temp folder
                convert(os.path.join(tmpdir, "upload.docx"), tmp_pdf_path)  # convert docx to pdf

                doc = fitz.open(tmp_pdf_path)  # open the converted pdf
                text = [p.get_text() for p in doc if p.get_text().strip()]  # extract text page by page
                if text:
                    return text  # return text if any

                return [pytesseract.image_to_string(img) for img in convert_from_path(tmp_pdf_path)]  # fallback ocr

        elif ext == ".pptx":
            prs = Presentation(file_path)  # open pptx file
            return [
                f"Slide {i+1}:\n" + "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
                ) for i, slide in enumerate(prs.slides)
            ]  # extract text from each slide

        elif ext in [".xls", ".xlsx"]:
            text_chunks = []  # init list to store text chunks
            dfs = pd.read_excel(file_path, sheet_name=None)  # read all sheets
            for sheet_name, df in dfs.items():  # iterate over each sheet
                df = df.fillna("").astype(str)  # fill empty cells and convert all to string
                for i, row in df.iterrows():  # iterate over rows
                    row_text = ", ".join(row.tolist())  # join row values into single line
                    if row_text.strip():
                        text_chunks.append(f"Sheet: {sheet_name} | Row: {i+1}\n{row_text}")  # format chunk
            return text_chunks  # return all collected chunks

    except Exception as e:
        print(f"❌ Error reading {file_path}: {e}")  # log any error

    return []  # return empty if nothing found

#RE RANKERS 
re_ranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2") # re ranker model: fast, accurate, local
def rerank_chunks(query, chunks, top_k=5):
    """
Re-rank retrieved document chunks based on semantic relevance to the query using a CrossEncoder.

Args:
    query (str): User input query to rank chunks against.
    chunks (List[Dict]): Retrieved chunk list from Pinecone with 'metadata.text' field.
    top_k (int): Number of top-ranked chunks to return.

Returns:
    List[Dict]: Top-k chunks, sorted by semantic relevance score (descending).
"""

    texts = [chunk["metadata"]["text"] for chunk in chunks]
    pairs = [(query, text) for text in texts]
    scores = re_ranker.predict(pairs)
    ranked = sorted(zip(chunks, scores), key=lambda x: x[1], reverse=True)
    return [chunk for chunk, _ in ranked[:top_k]]



def process_single_file(file):
    """Process and embed a single document into Pinecone and MongoDB."""
    try:
        filename = os.path.basename(file.name)
        #file_hash = calculate_file_hash(file.name)

        if docs_collection is not None:
            existing_doc = docs_collection.find_one({"filename": filename})
            if existing_doc:
                chunk_ids = existing_doc.get("chunk_ids", [])
                if chunk_ids:
                    print(f"🧹 Deleting {len(chunk_ids)} old chunks for {filename} from Pinecone")
                    index.delete(ids=chunk_ids, namespace=NAMESPACE)
                docs_collection.delete_one({"_id": existing_doc["_id"]})

        pages = extract_text(file.name)
        if not pages:
            return f"⚠️ No text found in {filename}"

        all_chunks = []
        for page_num, page_text in enumerate(pages, start=1):
            splits = splitter.split_text(page_text)
            for i, chunk in enumerate(splits):
                if not chunk.strip():
                    continue
                embedding = model.encode(chunk).tolist()
                meta = {
                    "source": filename,
                    "chunk": i + 1,
                    "text": chunk,
                    "page": page_num 
                }
                if "Slide" in page_text:
                    meta["slide"] = page_num
                elif "Sheet:" in page_text:
                    try:
                        meta["sheet"] = page_text.split("Sheet:")[1].split("|")[0].strip()
                        meta["row"] = page_num
                    except: pass
                all_chunks.append({
                    "id": f"{filename}-p{page_num}-c{i}",
                    "values": embedding,
                    "metadata": meta
                })
                # Force namespace creation if empty (Pinecone bug workaround)
        try:
            index.upsert(vectors=[{
                "id": "__init__",
                "values": [0.0] * 384,
                "metadata": {"source": "__init__"}
            }], namespace=NAMESPACE)
            index.delete(ids=["__init__"], namespace=NAMESPACE)
        except Exception as e:
            print(f"⚠️ Namespace init failed: {e}")


        index.upsert(vectors=all_chunks, namespace=NAMESPACE)

        docs_collection.insert_one({
            "filename": filename,
            #"hash": file_hash,
            "chunks": len(all_chunks),
            "pages": len(pages),
            "chunk_ids": [chunk["id"] for chunk in all_chunks],
            "uploaded_at": datetime.now(timezone.utc)
        })


        return f"✅ Uploaded {filename} ({len(all_chunks)} chunks)"
    except Exception as e:
        return f"❌ Upload failed: {e}"

def process_multiple_files(files):
    """Process multiple document uploads."""
    if not files:
        return "❌ No files uploaded."
    return "\n".join([process_single_file(f) for f in files])



def get_all_docs():
    """Retrieve list of all uploaded documents from MongoDB."""
    try:
        if docs_collection is None:
            return []
        return [doc["filename"] for doc in docs_collection.find()]
    except Exception as e:
        print(f"❌ Mongo fetch error: {e}")
        return []

def delete_docs(filenames):
    """Delete selected documents from MongoDB and Pinecone."""
    deleted = []
    for f in filenames:
        try:
            if docs_collection is not None:
                docs_collection.delete_many({"filename": f})
            index.delete(filter={"source": f}, namespace=NAMESPACE)
            deleted.append(f)
        except Exception as e:
            print(f"❌ Delete error: {e}")
    updated = get_all_docs()
    return f"🗑️ Deleted: {', '.join(deleted)}", gr.update(choices=updated, value=[])

def delete_all_docs():
    """Delete all document vectors from MongoDB and Pinecone, but keep namespace alive."""
    try:
        # If MongoDB is connected, delete all document metadata
        if docs_collection is not None:
            docs_collection.delete_many({})  # Remove all documents from the MongoDB collection

        # Query Pinecone to fetch existing vectors (up to 1000) in the namespace
        res = index.query(
            vector=[0.0]*384,          # Dummy vector (shape must match your embedding model, here 384-dim)
            top_k=1000,                # Max number of vectors to fetch
            include_metadata=True,     # Include metadata to identify '__init__'
            namespace=NAMESPACE        # Target namespace (keep it alive)
        )

        # Filter vector IDs to delete: skip dummy '__init__' vector
        ids_to_delete = [
            m["id"] for m in res.get("matches", [])               # Loop over returned matches
            if m["metadata"].get("source") != "__init__"          # Keep dummy vector for namespace preservation
        ]

        # If there are vectors to delete, delete them
        if ids_to_delete:
            index.delete(ids=ids_to_delete, namespace=NAMESPACE)  # Delete only actual document vectors

        # Return UI update: all documents deleted, empty doc list
        return "🗑️ All documents deleted (namespace preserved).", gr.update(choices=[], value=[])

    except Exception as e:
        # If an error occurs, return failure message and keep UI unchanged
        return f"❌ Delete all failed: {e}", gr.update()

def load_docs():
    """Reload document list for UI."""
    return gr.update(choices=get_all_docs(), value=[])

def get_top_chunks(query, embedder, index, top_k=5):
    """
    Retrieve the most relevant text chunks from the Pinecone vector index for a given user query.

    This function embeds the input query using the provided sentence embedding model,
    then performs a similarity search in the Pinecone vector database to find the top-k
    most relevant text chunks based on cosine similarity.

    Args:
        query (str): 
            The user's natural language question or search query.
            This will be converted into a dense vector for semantic search.
        
        embedder (SentenceTransformer): 
            The embedding model used to convert the query string into a numerical vector.
            Typically a pre-trained SentenceTransformer model like "all-MiniLM-L6-v2".

        index (pinecone.Index): 
            The Pinecone vector index instance to search against.
            This index contains previously embedded and stored document chunks.

        top_k (int, optional): 
            The number of most relevant document chunks to return. 
            Defaults to 5.

    Returns:
        list:
            A list of dictionary objects (Pinecone matches), each representing a document chunk.
            Each match contains:
                - "id": unique ID of the chunk
                - "score": similarity score with the query
                - "metadata": dictionary of associated metadata like filename, page, chunk number, text, etc.
    """
    vector = embedder.encode(query).tolist()  # Convert the query into a dense vector
    res = index.query(
        vector=vector,
        top_k=top_k,
        include_metadata=True,
        namespace=NAMESPACE
    )
    return res.get("matches", [])  # Return the list of matching chunks

def ask_llm(question, context, model_name="llama3-8b-8192"):
    
    """
    Send a question to a Large Language Model (LLM) using the Groq API, along with a specific context,
    and return a generated answer that is strictly based on the context.

    This function ensures that the LLM only uses the provided context to formulate an answer, 
    and provides detailed instructions to reference the source (filename and location) where 
    the information was found — useful in Retrieval-Augmented Generation (RAG) systems.

    Args:
        question (str): 
            The user's question to be answered by the LLM.
        
        context (str): 
            A string containing the contextual information retrieved from vector search (e.g., document chunks).
            The LLM is instructed to rely solely on this context to formulate an answer.

        model_name (str, optional): 
            The name of the LLM to use from the Groq API. 
            Defaults to "llama3-8b-8192".

    Returns:
        str:
            The LLM's generated answer as a string. 
            If an error occurs (e.g., network or API issue), it returns a string with the error message.

    Example:
        >>> context = "--- document.pdf (Page 3) ---\\nThis document talks about climate change..."
        >>> ask_llm("What does the document say about climate change?", context)
        "Climate change refers to long-term shifts in temperatures and weather patterns..."

    Notes:
        - The system prompt guides the LLM to:
            - Only use the provided context
            - Acknowledge missing information if not found in context
            - Include references such as page/slide/sheet numbers where relevant
    """

    try:
        response = groq_client.chat.completions.create(
            model=model_name,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an intelligent assistant. You must answer ALL user questions using only the provided context. "
                        "Give a detailed answer in accordance with the context."
                        "If a specific answer cannot be found in the context, clearly state: "
                        "'The context does not contain information about [that part]'. "
                        "Always mention the filename and page number/s where the information was found for PDF and DOCX. Mention under references section."
                        "Always mention the filename and slide number/s where the information was found for PPT and PPTX. Mention under references section."
                        "Always mention the filname and sheet/s and row number/s where the information was found for XLSX and XLS. Mention under references section."
                    )
                },
                {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}\nAnswer:"}
            ],
            max_tokens=384,
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"❌ LLM error: {e}"

def _run_rag_pipeline(q):
    """
Run the full Retrieval-Augmented Generation (RAG) pipeline:
1. Embeds query
2. Retrieves matching chunks from Pinecone
3. Re-ranks using a CrossEncoder
4. Passes the top-k to Groq LLM for answer generation

Args:
    q (str): User's input question.

Returns:
    str: Markdown-formatted answer and references.
"""

    if not q.strip():
        return "⚠️ Please enter a question."

    embedder, index = load_clients()

    matches = get_top_chunks(q, embedder, index, top_k=20)  # fetch more for better re-ranking
    matches = rerank_chunks(q, matches, top_k=5)  # keep top 5 most relevant

    if not matches:
        return "❌ No relevant context found."

    ctx, refs = "", []
    for m in matches:
        meta = m.get("metadata", {})
        fn = meta.get("source", "Unknown")
        txt = meta.get("text", "[No text found]")

        # Determine file type and format reference accordingly
        ext = fn.lower().split('.')[-1]
        location = None

        if ext in ["ppt", "pptx"]:
            slide = meta.get("slide")
            if slide: location = f"Slide {slide}"

        elif ext in ["xls", "xlsx"]:
            sheet = meta.get("sheet")
            row = meta.get("row")
            if sheet and row: location = f"Sheet {sheet}, Row {row}"

        elif ext in ["pdf", "docx"]:
            page = meta.get("page")
            if page: location = f"Page {page}"

        # Fallback only if we couldn't determine location
        if not location:
            location = "Location Unknown"

        ctx += f"\n--- {fn} ({location}) ---\n{txt}\n"
        refs.append(f"{fn} ({location})")

    ans = ask_llm(q, ctx)
    return f"### 🧠 Answer\n{ans}\n\n---\n\n### 📚 References\n" + "\n".join(f"- {r}" for r in refs)



# ---------------------- UI ----------------------

upload_tab = gr.Interface(  # file upload tab
    fn=process_multiple_files,
    inputs=gr.File(label="Upload documents", file_count="multiple"),
    outputs=gr.Textbox(label="Upload Status"),
    title="Upload Files",
    description="Upload PDF, DOCX, PPTX, or Excel files."
)

with gr.Blocks() as chat_tab:  # chat interface
    gr.Markdown("### 💬 Ask Questions About Your Files")
    question = gr.Textbox(label="Your Question", lines=2)
    btn = gr.Button("Submit")
    loading = gr.Markdown("⏳ Processing...", visible=False)
    output = gr.Markdown()
    btn.click(lambda q: (gr.update(visible=True), None), inputs=question, outputs=[loading, output])
    btn.click(lambda q: (gr.update(visible=False), _run_rag_pipeline(q)), inputs=question, outputs=[loading, output])

with gr.Blocks() as delete_tab:  # document deletion UI
    gr.Markdown("### 🗑️ Delete Documents")
    docs = gr.CheckboxGroup(label="Select documents to delete")
    delete_btn = gr.Button("Delete Selected")
    delete_all_btn = gr.Button("Delete All")
    refresh_btn = gr.Button("🔄 Refresh List")

    msg = gr.Textbox(label="Delete Status")
    delete_tab.load(load_docs, outputs=docs)
    delete_btn.click(delete_docs, inputs=docs, outputs=[msg, docs])
    delete_all_btn.click(delete_all_docs, outputs=[msg, docs])
    refresh_btn.click(load_docs, outputs=docs)

# Combine all tabs
app = gr.TabbedInterface([upload_tab, chat_tab, delete_tab], ["Upload", "Ask", "Delete"])

if __name__ == "__main__":
    app.launch(share=True)  # launch the Gradio app

