#python upload_to_pinecone.py


from dotenv import load_dotenv
load_dotenv()

import os
import time
from docx2pdf import convert 
import pdfplumber
import fitz
from docx import Document
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pinecone import Pinecone
from pptx import Presentation
import pandas as pd
import pytesseract
from pdf2image import convert_from_path


# --- Disable parallel tokenizer warning
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# --- Load env vars
pinecone_api_key = os.environ.get("PINECONE_API_KEY")
index_name = os.environ.get("PINECONE_INDEX", "rag-index")
pdf_folder = "data/"
batch_size = 50
start_batch = 0

print("🔑 Pinecone Key:", pinecone_api_key)
print("📌 Index:", index_name)

# --- Initialize Pinecone
pc = Pinecone(api_key=pinecone_api_key)
index_list = [i.name for i in pc.list_indexes()]

if index_name not in index_list:
    print(f"❌ Index '{index_name}' not found. Please create it in Pinecone first.")
    exit()

index = pc.Index(index_name)

# --- Load embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")

# --- Text splitter
splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
all_chunks = []

# --- Extract text from files
def extract_text(file_path):
    try:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            doc = fitz.open(file_path)
            text_pages = [page.get_text() for page in doc if page.get_text().strip()]
            if text_pages:
                return text_pages
            images = convert_from_path(file_path)
            return [pytesseract.image_to_string(img) for img in images if pytesseract.image_to_string(img).strip()]

        elif ext == ".docx":
            doc = Document(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]

        elif ext == ".pptx":
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
                if slide_text.strip():
                    slides_text.append(slide_text)
            return slides_text

        elif ext in [".xlsx", ".xls"]:
            df_list = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            text_chunks = []
            for sheet_name, df in df_list.items():
                text = df.astype(str).replace("nan", "").to_string(index=False, header=True)
                if text.strip():
                    text_chunks.append(f"Sheet: {sheet_name}\n{text}")
            return text_chunks

    except Exception as e:
        print(f"❌ error reading file {file_path}: {e}")

    return []

#uploading to pinecone
for filename in os.listdir(pdf_folder):
    file_path = os.path.join(pdf_folder, filename)

    # Determine extension
    ext = filename.lower().split('.')[-1]
    original_filename = filename  # Save original name for metadata

    # Handle PDF and DOCX (your original logic)
    if filename.endswith(".pdf") or filename.endswith(".docx"):
        if filename.endswith(".docx"):
            print(f"🌀 Converting {filename} to PDF...")
            try:
                converted_pdf_path = os.path.join(pdf_folder, filename.replace(".docx", ".pdf"))
                convert(file_path, converted_pdf_path)
                file_path = converted_pdf_path
                filename = os.path.basename(file_path)
            except Exception as e:
                print(f"⚠️ Failed to convert {filename} to PDF: {e}")
                continue
        print(f"📄 Processing {original_filename}...")
        page_texts = extract_text(file_path)

    # Handle PPT or PPTX
    elif filename.endswith(".pptx"):
        print(f"📊 Processing PowerPoint: {original_filename}...")
        page_texts = extract_text(file_path)

    # Handle XLS or XLSX
    elif filename.endswith(".xls") or filename.endswith(".xlsx"):
        print(f"📈 Processing Excel: {original_filename}...")
        page_texts = extract_text(file_path)

    else:
        continue  # Skip unsupported files

    # Process extracted pages/texts
    for page_num, page_text in enumerate(page_texts, start=1):
        if not page_text.strip():
            continue
        splits = splitter.split_text(page_text)
        for i, chunk in enumerate(splits):
            chunk_id = f"{original_filename}-p{page_num}-c{i}"  # Preserve original filename
            all_chunks.append({
                "id": chunk_id,
                "text": chunk,
                "metadata": {
                    "source": original_filename,
                    "text": i + 1,
                    "page": page_num
                }
            })

# --- Upload to Pinecone
print(f"📦 Total chunks to upload: {len(all_chunks)}")

for i in range(start_batch * batch_size, len(all_chunks), batch_size):
    batch = all_chunks[i:i + batch_size]
    print(f"🚀 Uploading batch {i // batch_size + 1} of {(len(all_chunks) - 1) // batch_size + 1}")
    vectors = []

    for item in batch:
        try:
            embedding = model.encode(item["text"]).tolist()
            vectors.append({
                "id": item["id"],
                "values": embedding,
                "metadata": {
                    **item["metadata"],
                    "text": item["text"]  # put actual content into metadata
                }
            })

        except Exception as e:
            print(f"⚠️ Failed to embed {item['id']}: {e}")
            continue

    if vectors:
        try:
            index.upsert(vectors=vectors)
            print(f"✅ Uploaded {len(vectors)} vectors.")
        except Exception as e:
            print(f"❌ Failed to upload batch: {e}")
            break
    else:
        print("⚠️ No valid vectors in this batch.")
    time.sleep(1)

print("✅ All Done!")
