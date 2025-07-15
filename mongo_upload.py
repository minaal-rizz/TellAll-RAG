import os  # to interact with the operating system
import fitz  # pymupdf - for reading pdf content
from docx import Document  # for reading .docx files
from sentence_transformers import SentenceTransformer  # for generating text embeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter  # to split text into chunks
from pymongo import MongoClient  # to interact with mongodb
from pinecone import Pinecone  # to interact with pinecone vector db
import gradio as gr  # for creating web ui
from dotenv import load_dotenv  # to load env vars from .env file
from pdf2image import convert_from_path  # convert pdf pages to images
import pytesseract  # perform ocr on images
from pptx import Presentation
import pandas as pd


# --- load environment variables
load_dotenv()  # load variables from .env file
pinecone_key = os.getenv("PINECONE_API_KEY")  # get pinecone api key
pinecone_index = os.getenv("PINECONE_INDEX")  # get pinecone index name
mongo_uri = os.getenv("MONGO_URI")  # get mongodb uri

# --- initialize clients
model = SentenceTransformer("all-MiniLM-L6-v2")  # load sentence transformer model
splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)  # define chunking strategy
pc = Pinecone(api_key=pinecone_key)  # initialize pinecone client
index = pc.Index(pinecone_index)  # access the specified pinecone index
client = MongoClient(mongo_uri)  # connect to mongodb
db = client["tellall"]  # select mongodb database
docs_collection = db["documents"]  # use the documents collection

# --- extract text from pdf or docx

def extract_text(file_path):
    try:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            doc = fitz.open(file_path)
            text_pages = [page.get_text() for page in doc if page.get_text().strip()]
            if text_pages:
                return text_pages
            images = convert_from_path(file_path)
            return [pytesseract.image_to_string(img) for img in images if pytesseract.image_to_string(img).strip()]

        elif ext == ".docx":
            doc = Document(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]

        elif ext == ".pptx":
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
                if slide_text.strip():
                    slides_text.append(slide_text)
            return slides_text

        elif ext in [".xlsx", ".xls"]:
            df_list = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            text_chunks = []
            for sheet_name, df in df_list.items():
                text = df.astype(str).replace("nan", "").to_string(index=False, header=True)
                if text.strip():
                    text_chunks.append(f"Sheet: {sheet_name}\n{text}")
            return text_chunks

    except Exception as e:
        print(f"❌ error reading file {file_path}: {e}")

    return []
# --- process and upload a single file
def process_single_file(file):  # handle one file
    filename = os.path.basename(file.name)  # extract filename
    file_path = file.name  # get file path
    text_pages = extract_text(file_path)  # extract text content from file

    if not text_pages:
        return f"⚠️ no valid text found in {filename}"  # skip if no content

    all_chunks = []  # to hold all processed text chunks

    for page_num, page_text in enumerate(text_pages, start=1):  # go through pages
        splits = splitter.split_text(page_text)  # split each page into chunks
        for i, chunk in enumerate(splits):  # go through each chunk
            if not chunk.strip():
                continue  # skip empty chunks
            try:
                embedding = model.encode(chunk).tolist()  # generate embedding
                chunk_id = f"{filename}-p{page_num}-c{i}"  # unique chunk id
                all_chunks.append({  # collect chunk data
                    "id": chunk_id,
                    "text": chunk,
                    "metadata": {
                        "source": filename,
                        
                    },
                    "embedding": embedding
                })
            except Exception as e:
                print(f"❌ embedding error in {filename} (p{page_num} c{i}): {e}")  # log if embedding fails

    vectors = [{  # prepare list for pinecone upsert
        "id": item["id"],
        "values": item["embedding"],
        "metadata": item["metadata"]
    } for item in all_chunks]

    try:
        index.upsert(vectors=vectors)  # upload embeddings to pinecone
        docs_collection.insert_one({  # save document metadata to mongodb
            "filename": filename,
            "chunks": len(all_chunks),
            "pages": len(text_pages),
            
        })
        return f"✅ uploaded {filename} ({len(all_chunks)} chunks)"  # return success
    except Exception as e:
        return f"❌ upload failed for {filename}: {e}"  # return error

# --- upload and process multiple files
def process_multiple_files(file_list):  # handle batch uploads
    if not file_list:
        return "❌ no files uploaded."  # return if nothing uploaded
    results = []  # to hold output messages
    for file in file_list:
        try:
            result = process_single_file(file)  # process one file
            results.append(result)  # add result
        except Exception as e:
            results.append(f"❌ error with {file.name}: {e}")  # handle exceptions
    return "\n".join(results)  # return output

# --- sync pinecone namespaces to mongodb if missing

'''def sync_pinecone_to_mongo():  # sync pinecone to mongodb
    try:
        mongo_filenames = {doc["filename"] for doc in docs_collection.find({}, {"filename": 1})}  # get all mongo filenames
        pinecone_namespaces = index.describe_index_stats()["namespaces"].keys()  # get pinecone namespaces
        missing = set(pinecone_namespaces) - mongo_filenames  # find missing files

        for fname in missing:
            docs_collection.insert_one({  # insert missing file metadata
                "filename": fname,
                "synced_from_pinecone": True,
             
            })
            print(f"✅ synced {fname} to mongodb")  # log success

        if not missing:
            print("✅ mongodb already up to date with pinecone.")  # log if already synced
    except Exception as e:
        print(f"❌ failed syncing pinecone -> mongodb: {e}")'''


def sync_pinecone_to_mongo():
    try:
        mongo_filenames = {doc["filename"] for doc in docs_collection.find({}, {"filename": 1})}  # filenames in mongo

        # get vector metadata stats
        stats = index.describe_index_stats()
        all_ids = []

        # extract vector metadata from all entries (no namespace)
        metadata_entries = stats.get("namespaces", {}).get("", {}).get("metadata", [])

        filenames_in_pinecone = set()

        for entry in metadata_entries:
            source = entry.get("source")
            if source:
                filenames_in_pinecone.add(source)

        # find missing files in MongoDB
        missing = filenames_in_pinecone - mongo_filenames

        for fname in missing:
            docs_collection.insert_one({
                "filename": fname,
                "synced_from_pinecone": True,
               
            })
            print(f"✅ synced {fname} to mongodb")

        if not missing:
            print("✅ mongodb already up to date with pinecone.")

    except Exception as e:
        print(f"❌ failed syncing pinecone -> mongodb: {e}")

# --- gradio ui
upload_ui = gr.Interface(
    fn=process_multiple_files,  # main function for gradio
    inputs=gr.File(label="📄 upload pdfs/docxs", file_types=[".pdf", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"], file_count="multiple"),  # file input
    outputs=gr.Textbox(label="📦 upload status"),  # output text box
    title="🧠 TellAll: RAG Chatbot",  # gradio title
    description="upload and store documents in pinecone and mongodb."  # gradio description
)

# --- launch
if __name__ == "__main__":
    sync_pinecone_to_mongo()  # sync pinecone docs before ui starts
    upload_ui.launch()  # launch gradio app

